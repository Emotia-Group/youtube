"""Extracción de texto de documentos (guiones en PDF, Word, PowerPoint,
Excel, txt/md) para la fase de ingesta."""
from __future__ import annotations

from pathlib import Path


def _need(package: str, pip_name: str):
    raise RuntimeError(
        f"Para leer este tipo de archivo instala '{pip_name}' "
        f"(pip install {pip_name}).")


def extract_text(path: Path) -> str:
    """Devuelve el texto plano de un documento según su extensión."""
    ext = path.suffix.lower()

    if ext in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="replace")

    if ext == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            _need("pypdf", "pypdf")
        reader = PdfReader(str(path))
        return "\n\n".join((page.extract_text() or "") for page in reader.pages)

    if ext == ".docx":
        try:
            import docx
        except ImportError:
            _need("docx", "python-docx")
        document = docx.Document(str(path))
        return "\n".join(p.text for p in document.paragraphs)

    if ext == ".doc":
        raise RuntimeError(
            "El formato .doc antiguo no es compatible — guarda el archivo "
            "como .docx (Word moderno) o .pdf y vuelve a subirlo.")

    if ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            _need("pptx", "python-pptx")
        pres = Presentation(str(path))
        slides = []
        for i, slide in enumerate(pres.slides, 1):
            texts = [shape.text for shape in slide.shapes
                     if getattr(shape, "has_text_frame", False) and shape.text]
            if texts:
                slides.append(f"[Diapositiva {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)

    if ext in (".xlsx", ".xls"):
        try:
            from openpyxl import load_workbook
        except ImportError:
            _need("openpyxl", "openpyxl")
        wb = load_workbook(str(path), read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
        return "\n".join(rows)

    raise RuntimeError(f"No sé extraer texto de archivos {ext}.")
